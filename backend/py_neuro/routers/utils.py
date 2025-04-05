# Standard library imports
import os
import json
import logging
import uuid # Added for generating section IDs
from typing import List

# FastAPI imports
from fastapi import HTTPException

# PDF processing
import fitz  # PyMuPDF

# Document processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document

# Image processing
from PIL import Image, UnidentifiedImageError

# Google AI imports
import google.generativeai as genai

# Local imports
from .. import models # Relative import
from pydantic import ValidationError # Import for specific error checking
# from utils.image_processing import pre_filter_image
from .. import config # Import the whole module

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

TEMP_UPLOAD_DIR = "temp_uploads"
IMAGE_OUTPUT_DIR_BASE = os.path.join(TEMP_UPLOAD_DIR, "images")
MIN_IMAGE_WIDTH = 50  # Pixels - adjust as needed for pre-filtering
MIN_IMAGE_HEIGHT = 50 # Pixels - adjust as needed for pre-filtering

def extract_text_and_images_pdf(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo], List[models.ExtractedImageInfo]]:
    """Extracts text and images from PDF, performs basic image filtering."""
    text = ""
    all_images_info: List[models.ExtractedImageInfo] = []
    filtered_images_info: List[models.ExtractedImageInfo] = []
    try:
        doc = fitz.open(file_path)
        img_counter = 0
        for page_num, page in enumerate(doc):
            text += page.get_text() + "\n"
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_filename = f"img_{page_num + 1}_{img_counter}.{image_ext}"
                    image_save_path = os.path.join(request_image_dir, image_filename)

                    with open(image_save_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    image_info = models.ExtractedImageInfo(filename=image_filename, page_number=page_num + 1)
                    all_images_info.append(image_info) # Add to all list first

                    # Pre-filter the saved image
                    if pre_filter_image(image_save_path):
                         filtered_images_info.append(image_info) # Add to filtered list if it passes
                    else:
                         # Optionally delete the filtered-out image file immediately
                         try: os.remove(image_save_path)
                         except OSError: pass
                    img_counter += 1
                except Exception as img_extract_err:
                     logger.warning(f"Could not extract image ref {xref} on page {page_num+1}: {img_extract_err}")

        logger.info(f"Extracted text, found {len(all_images_info)} total images, kept {len(filtered_images_info)} after filtering from PDF: {file_path}")
        doc.close()
        return text, all_images_info, filtered_images_info
    except Exception as e:
        logger.error(f"Error extracting data from PDF {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PDF file: {e}")

def extract_text_and_images_pptx(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo], List[models.ExtractedImageInfo]]:
    """Extracts text and images from PPTX, performs basic image filtering."""
    text = ""
    all_images_info: List[models.ExtractedImageInfo] = []
    filtered_images_info: List[models.ExtractedImageInfo] = []
    try:
        prs = Presentation(file_path)
        img_counter = 0
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext.lower()
                        image_filename = f"img_slide_{slide_num + 1}_{img_counter}.{image_ext}"
                        image_save_path = os.path.join(request_image_dir, image_filename)

                        with open(image_save_path, "wb") as img_file:
                            img_file.write(image_bytes)

                        image_info = models.ExtractedImageInfo(filename=image_filename, page_number=slide_num + 1)
                        all_images_info.append(image_info)

                        if pre_filter_image(image_save_path):
                            filtered_images_info.append(image_info)
                        else:
                            try: os.remove(image_save_path)
                            except OSError: pass
                        img_counter += 1
                    except Exception as img_extract_err:
                        logger.warning(f"Could not extract image shape on slide {slide_num+1}: {img_extract_err}")

        logger.info(f"Extracted text, found {len(all_images_info)} total images, kept {len(filtered_images_info)} after filtering from PPTX: {file_path}")
        return text, all_images_info, filtered_images_info
    except Exception as e:
        logger.error(f"Error extracting data from PPTX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX file: {e}")

def extract_text_docx(file_path: str) -> str:
    """Extracts text from a DOCX file. Image extraction is complex and omitted for now."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        logger.info(f"Extracted text from DOCX: {file_path}. Image extraction skipped.")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process DOCX file: {e}")

def extract_text_txt(file_path: str) -> str:
    """Extracts text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        logger.info(f"Extracted text from TXT: {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from TXT {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process TXT file: {e}")


# Modified extraction functions to return BOTH all extracted and filtered images
async def generate_hierarchical_study_guide(text_content: str, images_info: List[models.ExtractedImageInfo], request_image_dir: str) -> List[models.StudyGuideSection]:
    """Uses Gemini multimodal capabilities to generate hierarchical sections, explanations, and image associations."""
    if not config.gemini_model: # Access via config module
        raise HTTPException(status_code=500, detail="Gemini model not configured. Check API key.")

    prepared_images = []
    image_filenames_list = [] # Only filenames of images being sent (already pre-filtered)
    for img_info in images_info: # images_info now contains only pre-filtered images
        try:
            img_path = os.path.join(request_image_dir, img_info.filename)
            pil_image = Image.open(img_path)
            if pil_image.mode != 'RGB':
                 pil_image = pil_image.convert('RGB')
            prepared_images.append(pil_image)
            image_filenames_list.append(img_info.filename)
        except Exception as e:
            logger.warning(f"Could not load or prepare image {img_info.filename} for Gemini: {e}")

    # --- Prompt v4: Hierarchical Structure, Relevance Filter, Markdown Cues ---
    prompt_parts = [
        """You are an expert academic assistant tasked with creating an enhanced, hierarchical study guide from a document's text and associated images. Your goal is to break down the content into logical main sections, further divide those into subsections, explain the concepts within each subsection clearly like a tutor, and integrate explanations of relevant images directly into the text.

**Input:**
You will receive the full text content of the document and a list of image filenames corresponding to potentially relevant images extracted from the document (provided alongside the text).

**Tasks:**

1.  **Analyze Text:** Read through the entire provided text content.
2.  **Identify Main Sections:** Divide the text into logical **main sections** based on the primary topics or chapters.
3.  **For Each Main Section:**
    *   Generate a concise `section_title`.
    *   Generate a brief `section_overview_description` (1-2 sentences) summarizing the section's goal or content.
    *   Identify logical **subsections** within this main section's text.
    *   Generate a list of `subsection_titles` for these identified subsections.
    *   For **each subsection**:
        *   Generate a detailed `explanation` (aim for **at least 5-8 sentences**, or more if needed for clarity) that:
            *   Thoroughly clarifies the key concepts discussed within that subsection's text.
            *   Defines important jargon or technical terms.
            *   Simplifies complex ideas using analogies or simpler language where appropriate.
            *   Provides necessary background context if helpful for understanding.
            *   **Use Markdown syntax within this explanation text for formatting:** Use asterisks (`*` or `-`) for bullet points, double asterisks (`**text**`) for bolding, backticks (`) for inline code, and LaTeX delimiters (`$...$` or `$$...$$`) for mathematical formulas. Use Markdown table syntax if needed.
            *   **Crucially:** Review the provided **filtered** images (corresponding to the filenames listed below) and determine which images are directly relevant to the content being explained in **this specific subsection**. **Only associate images that directly illustrate or clarify the concepts. Ignore purely decorative images, logos, or images that appear blank or lack meaningful content relevant to the text.** If relevant images are found, **explicitly reference them within this explanation text** (e.g., "Figure [filename] illustrates...", "As seen in the diagram [filename], the process involves...", "The results in [filename] show..."). Explain what the referenced image depicts and how it connects to the concepts being explained.
        *   Generate a list `associated_image_filenames` containing the filenames of the images you identified as relevant *and referenced* in the explanation for this subsection. If no images are relevant or referenced, provide an empty list `[]`.
    *   **Focus:** Prioritize deep explanation and enhancing understanding over mere summarization. Do not simply rephrase the original text; add value through clarification and context.
4.  **Format Output:** Structure your entire response as a single JSON list of objects. Each object represents a **main section** and must contain the following keys:
    *   `section_title` (string): The title for the main section.
    *   `section_overview_description` (string): The brief overview description.
    *   `subsection_titles` (list of strings): The titles of the subsections within this section.
    *   `subsections` (list of objects): A list where each object represents a subsection and contains:
        *   `subsection_title` (string): The title of the subsection.
        *   `explanation` (string): The detailed explanatory text for the subsection, including Markdown and image references.
        *   `associated_image_filenames` (list of strings): Filenames of images relevant to *this subsection*.

**Output Constraints:**
*   Your response MUST be only the JSON list of main section objects, starting with `[` and ending with `]`.
*   Do NOT include any introductory text, concluding remarks, or markdown formatting (like ```json) around the JSON output.
*   Ensure the generated JSON is valid and strictly adheres to the nested structure described above.

**Example JSON Object (representing one main section):**
```json
{
  "section_title": "Understanding PinSage Architecture",
  "section_overview_description": "This section details the core components and workflow of the PinSage algorithm.",
  "subsection_titles": ["Localized Convolutions", "Minibatch Construction"],
  "subsections": [
    {
      "subsection_title": "Localized Convolutions",
      "explanation": "PinSage employs localized convolutions, focusing on a node's local neighborhood rather than the entire graph. This is visualized in **Figure img_6_2.png**. The process involves sampling neighbors via random walks and aggregating features, weighted by importance using techniques like importance pooling. This avoids the computational bottleneck of processing the full graph Laplacian.",
      "associated_image_filenames": ["img_6_2.png"]
    },
    {
      "subsection_title": "Minibatch Construction",
      "explanation": "To train efficiently on massive graphs, PinSage uses a producer-consumer minibatch approach, conceptually illustrated in **diagram img_7_3.png**. This allows parallel preparation of computation graphs on the CPU while the GPU performs model training, maximizing hardware utilization and enabling large batch sizes.",
      "associated_image_filenames": ["img_7_3.png"]
    }
  ]
}
```

**Input Data:**

**Document Text:**
---
""",
        text_content[:30000], # Limit text length
        "\n---\n\n**Available (Pre-filtered) Image Filenames (corresponding to the images provided):**\n",
        json.dumps(image_filenames_list),
        "\n\n**JSON Output:**\n"
    ]

    for img in prepared_images:
        prompt_parts.append(img)

    # Use explicit JSON mode with gemini-1.5-pro
    generation_config = genai.types.GenerationConfig(response_mime_type="application/json")

    try:
        logger.info(f"Sending hierarchical multimodal request to Gemini with {len(prepared_images)} images...")
        response = await config.gemini_model.generate_content_async( # Access via config module
            prompt_parts,
            generation_config=generation_config
        )

        cleaned_response_text = response.text.strip().strip('```json').strip('```').strip()
        logger.info(f"Received raw response from Gemini: {response.text}")
        logger.info(f"Cleaned response text for JSON parsing: {cleaned_response_text}")

        # --- MODIFICATION START ---
        # 1. Find the last closing bracket and truncate the string
        last_bracket_index = cleaned_response_text.rfind(']')
        if last_bracket_index != -1:
            json_to_parse = cleaned_response_text[:last_bracket_index + 1]
            logger.info(f"Truncated response for parsing: {json_to_parse}")
        else:
            # If no closing bracket is found, attempt to parse the original cleaned text
            # but log a warning as it's likely malformed.
            logger.warning("Could not find closing bracket ']' in cleaned AI response. Attempting to parse anyway.")
            json_to_parse = cleaned_response_text

        # 2. Parse the potentially truncated JSON string into a Python list of dicts
        raw_sections_list = json.loads(json_to_parse)

        # 3. Add section_id to each section dictionary
        for section_dict in raw_sections_list:
            # Ensure it's actually a dictionary before adding the key
            if isinstance(section_dict, dict):
                 section_dict["section_id"] = str(uuid.uuid4())
            else:
                 # Handle cases where the AI might return non-dict items in the list
                 logger.warning(f"Skipping non-dictionary item found in AI response list: {section_dict}")


        # 3. Prepare the data for Pydantic validation (wrap in the expected structure)
        data_for_validation = {"study_guide": raw_sections_list}
        # --- MODIFICATION END ---

        # Use Pydantic to parse and validate the list of sections
        # Need a temporary wrapper for the top-level list
        class TempWrapper(models.BaseModel):
             # Use the correct hierarchical model name: StudyGuideSection
             study_guide: List[models.StudyGuideSection]

        # Validate the modified data
        parsed_wrapper = TempWrapper.model_validate(data_for_validation) # Use model_validate for dict
        parsed_sections = parsed_wrapper.study_guide

        logger.info(f"Successfully parsed Gemini response into {len(parsed_sections)} main sections.")
        return parsed_sections

    except json.JSONDecodeError as json_err: # Add specific catch for JSON parsing
        logger.error(f"Failed to parse Gemini response as JSON: {json_err}")
        logger.error(f"Cleaned text that failed JSON parsing: {cleaned_response_text if 'cleaned_response_text' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"AI response was not valid JSON: {json_err}")
    except ValidationError as val_err: # Catch Pydantic validation errors specifically
        logger.error(f"Pydantic validation failed after adding section_id: {val_err}")
        logger.error(f"Data failing validation: {data_for_validation if 'data_for_validation' in locals() else 'N/A'}")
        logger.error(f"Gemini raw response was: {response.text if 'response' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"AI response structure mismatch after processing: {val_err}")
    except Exception as e: # General catch-all
        logger.error(f"Unexpected error generating/processing study guide: {e}")
        # Log relevant context if available
        logger.error(f"Cleaned text that failed: {cleaned_response_text if 'cleaned_response_text' in locals() else 'N/A'}")
        logger.error(f"Gemini raw response was: {response.text if 'response' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"Unexpected error processing AI response: {e}")

def pre_filter_image(image_path: str) -> bool:
    """Basic pre-filtering for images based on dimensions."""
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                logger.info(f"Filtering out image (too small): {image_path}")
                return False
            # Add other checks here if needed (e.g., color variance)
        return True
    except UnidentifiedImageError:
        logger.warning(f"Cannot identify image file, filtering out: {image_path}")
        return False
    except Exception as e:
        logger.warning(f"Error during image pre-filtering for {image_path}: {e}")
        return False # Filter out if error occurs
