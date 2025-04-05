import os
import shutil
import uuid
import json
from fastapi import APIRouter, UploadFile, File, HTTPException, Depends
from typing import List, Dict, Any
import google.generativeai as genai
# Remove dotenv import, config handles it
# from dotenv import load_dotenv
import logging
from PIL import Image, UnidentifiedImageError # For image handling & filtering
import io

# Text/Image extraction libraries
import fitz # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from bson import ObjectId # Import ObjectId

# Use absolute import
from py_neuro import models
from py_neuro.config import settings # Import settings
from py_neuro.database import get_database, get_gridfs_bucket # Import DB and GridFS dependencies
from py_neuro import crud # Import crud functions
from motor.motor_asyncio import AsyncIOMotorDatabase, AsyncIOMotorGridFSBucket # Import DB and GridFS types

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()

# --- Configuration ---
# load_dotenv() # Removed - handled by config.py
# GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") # Removed - use settings
TEMP_UPLOAD_DIR = "temp_uploads"
IMAGE_OUTPUT_DIR_BASE = os.path.join(TEMP_UPLOAD_DIR, "images")
MIN_IMAGE_WIDTH = 50  # Pixels - adjust as needed for pre-filtering
MIN_IMAGE_HEIGHT = 50 # Pixels - adjust as needed for pre-filtering

# Configure Gemini
if not settings.GEMINI_API_KEY or settings.GEMINI_API_KEY == "YOUR_GEMINI_API_KEY_HERE": # Check settings
    logger.warning("GEMINI_API_KEY not found in settings or is default. Study guide generation will fail.")
    gemini_model = None
else:
    try:
        genai.configure(api_key=settings.GEMINI_API_KEY) # Use settings
        # Using gemini-1.5-pro which supports multimodal and JSON output mode
        gemini_model = genai.GenerativeModel('gemini-2.0-flash')
        logger.info("Gemini configured successfully with gemini-2.0-flas.")
    except Exception as e:
        logger.error(f"Failed to configure Gemini: {e}")
        gemini_model = None

# Ensure temp directories exist
os.makedirs(TEMP_UPLOAD_DIR, exist_ok=True)
os.makedirs(IMAGE_OUTPUT_DIR_BASE, exist_ok=True)

# --- Helper Functions ---

def pre_filter_image(image_path: str) -> bool:
    """Basic pre-filtering for images based on dimensions."""
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                logger.info(f"Filtering out image (too small): {image_path}")
                return False
            # Add other checks here if needed (e.g., color variance)
        return True
    except UnidentifiedImageError:
        logger.warning(f"Cannot identify image file, filtering out: {image_path}")
        return False
    except Exception as e:
        logger.warning(f"Error during image pre-filtering for {image_path}: {e}")
        return False # Filter out if error occurs

# Modified extraction functions to return BOTH all extracted and filtered images
def extract_text_and_images_pdf(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo], List[models.ExtractedImageInfo]]:
    """Extracts text and images from PDF, performs basic image filtering."""
    text = ""
    all_images_info: List[models.ExtractedImageInfo] = []
    filtered_images_info: List[models.ExtractedImageInfo] = []
    try:
        doc = fitz.open(file_path)
        img_counter = 0
        for page_num, page in enumerate(doc):
            text += page.get_text() + "\n"
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_filename = f"img_{page_num + 1}_{img_counter}.{image_ext}"
                    image_save_path = os.path.join(request_image_dir, image_filename)

                    with open(image_save_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    image_info = models.ExtractedImageInfo(filename=image_filename, page_number=page_num + 1)
                    all_images_info.append(image_info) # Add to all list first

                    # Pre-filter the saved image
                    if pre_filter_image(image_save_path):
                         filtered_images_info.append(image_info) # Add to filtered list if it passes
                    else:
                         # Optionally delete the filtered-out image file immediately
                         try: os.remove(image_save_path)
                         except OSError: pass
                    img_counter += 1
                except Exception as img_extract_err:
                     logger.warning(f"Could not extract image ref {xref} on page {page_num+1}: {img_extract_err}")

        logger.info(f"Extracted text, found {len(all_images_info)} total images, kept {len(filtered_images_info)} after filtering from PDF: {file_path}")
        doc.close()
        return text, all_images_info, filtered_images_info
    except Exception as e:
        logger.error(f"Error extracting data from PDF {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PDF file: {e}")

def extract_text_and_images_pptx(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo], List[models.ExtractedImageInfo]]:
    """Extracts text and images from PPTX, performs basic image filtering."""
    text = ""
    all_images_info: List[models.ExtractedImageInfo] = []
    filtered_images_info: List[models.ExtractedImageInfo] = []
    try:
        prs = Presentation(file_path)
        img_counter = 0
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext.lower()
                        image_filename = f"img_slide_{slide_num + 1}_{img_counter}.{image_ext}"
                        image_save_path = os.path.join(request_image_dir, image_filename)

                        with open(image_save_path, "wb") as img_file:
                            img_file.write(image_bytes)

                        image_info = models.ExtractedImageInfo(filename=image_filename, page_number=slide_num + 1)
                        all_images_info.append(image_info)

                        if pre_filter_image(image_save_path):
                            filtered_images_info.append(image_info)
                        else:
                            try: os.remove(image_save_path)
                            except OSError: pass
                        img_counter += 1
                    except Exception as img_extract_err:
                        logger.warning(f"Could not extract image shape on slide {slide_num+1}: {img_extract_err}")

        logger.info(f"Extracted text, found {len(all_images_info)} total images, kept {len(filtered_images_info)} after filtering from PPTX: {file_path}")
        return text, all_images_info, filtered_images_info
    except Exception as e:
        logger.error(f"Error extracting data from PPTX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX file: {e}")

def extract_text_docx(file_path: str) -> str:
    """Extracts text from a DOCX file. Image extraction is complex and omitted for now."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        logger.info(f"Extracted text from DOCX: {file_path}. Image extraction skipped.")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process DOCX file: {e}")

def extract_text_txt(file_path: str) -> str:
    """Extracts text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        logger.info(f"Extracted text from TXT: {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from TXT {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process TXT file: {e}")

# --- Updated Gemini Interaction ---
# Renamed function and updated return type hint
async def generate_hierarchical_study_guide(text_content: str, images_info: List[models.ExtractedImageInfo], request_image_dir: str) -> List[models.StudyGuideSection]:
    """Uses Gemini multimodal capabilities to generate hierarchical sections, explanations, and image associations."""
    if not gemini_model:
        raise HTTPException(status_code=500, detail="Gemini model not configured. Check API key.")

    prepared_images = []
    image_filenames_list = [] # Only filenames of images being sent (already pre-filtered)
    for img_info in images_info: # images_info now contains only pre-filtered images
        try:
            img_path = os.path.join(request_image_dir, img_info.filename)
            pil_image = Image.open(img_path)
            if pil_image.mode != 'RGB':
                 pil_image = pil_image.convert('RGB')
            prepared_images.append(pil_image)
            image_filenames_list.append(img_info.filename)
        except Exception as e:
            logger.warning(f"Could not load or prepare image {img_info.filename} for Gemini: {e}")

    # --- Prompt v4: Hierarchical Structure, Relevance Filter, Markdown Cues ---
    prompt_parts = [
        """You are an expert academic assistant tasked with creating an enhanced, hierarchical study guide from a document's text and associated images. Your goal is to break down the content into logical main sections, further divide those into subsections, explain the concepts within each subsection clearly like a tutor, and integrate explanations of relevant images directly into the text.

**Input:**
You will receive the full text content of the document and a list of image filenames corresponding to potentially relevant images extracted from the document (provided alongside the text).

**Tasks:**

1.  **Analyze Text:** Read through the entire provided text content.
2.  **Identify Main Sections:** Divide the text into logical **main sections** based on the primary topics or chapters.
3.  **For Each Main Section:**
    *   Generate a concise `section_title`.
    *   Generate a brief `section_overview_description` (1-2 sentences) summarizing the section's goal or content.
    *   Identify logical **subsections** within this main section's text.
    *   Generate a list of `subsection_titles` for these identified subsections.
    *   For **each subsection**:
        *   Generate a detailed `explanation` (aim for **at least 5-8 sentences**, or more if needed for clarity) that:
            *   Thoroughly clarifies the key concepts discussed within that subsection's text.
            *   Defines important jargon or technical terms.
            *   Simplifies complex ideas using analogies or simpler language where appropriate.
            *   Provides necessary background context if helpful for understanding.
            *   **Use Markdown syntax within this explanation text for formatting:** Use asterisks (`*` or `-`) for bullet points, double asterisks (`**text**`) for bolding, backticks (`) for inline code, and LaTeX delimiters (`$...$` or `$$...$$`) for mathematical formulas. Use Markdown table syntax if needed.
            *   **Crucially:** Review the provided **filtered** images (corresponding to the filenames listed below) and determine which images are directly relevant to the content being explained in **this specific subsection**. **Only associate images that directly illustrate or clarify the concepts. Ignore purely decorative images, logos, or images that appear blank or lack meaningful content relevant to the text.** If relevant images are found, **explicitly reference them within this explanation text** (e.g., "Figure [filename] illustrates...", "As seen in the diagram [filename], the process involves...", "The results in [filename] show..."). Explain what the referenced image depicts and how it connects to the concepts being explained.
        *   Generate a list `associated_image_filenames` containing the filenames of the images you identified as relevant *and referenced* in the explanation for this subsection. If no images are relevant or referenced, provide an empty list `[]`.
    *   **Focus:** Prioritize deep explanation and enhancing understanding over mere summarization. Do not simply rephrase the original text; add value through clarification and context.
4.  **Format Output:** Structure your entire response as a single JSON list of objects. Each object represents a **main section** and must contain the following keys:
    *   `section_title` (string): The title for the main section.
    *   `section_overview_description` (string): The brief overview description.
    *   `subsection_titles` (list of strings): The titles of the subsections within this section.
    *   `subsections` (list of objects): A list where each object represents a subsection and contains:
        *   `subsection_title` (string): The title of the subsection.
        *   `explanation` (string): The detailed explanatory text for the subsection, including Markdown and image references.
        *   `associated_image_filenames` (list of strings): Filenames of images relevant to *this subsection*.

**Output Constraints:**
*   Your response MUST be only the JSON list of main section objects, starting with `[` and ending with `]`.
*   Do NOT include any introductory text, concluding remarks, or markdown formatting (like ```json) around the JSON output.
*   Ensure the generated JSON is valid and strictly adheres to the nested structure described above.

**Example JSON Object (representing one main section):**
```json
{
  "section_title": "Understanding PinSage Architecture",
  "section_overview_description": "This section details the core components and workflow of the PinSage algorithm.",
  "subsection_titles": ["Localized Convolutions", "Minibatch Construction"],
  "subsections": [
    {
      "subsection_title": "Localized Convolutions",
      "explanation": "PinSage employs localized convolutions, focusing on a node's local neighborhood rather than the entire graph. This is visualized in **Figure img_6_2.png**. The process involves sampling neighbors via random walks and aggregating features, weighted by importance using techniques like importance pooling. This avoids the computational bottleneck of processing the full graph Laplacian.",
      "associated_image_filenames": ["img_6_2.png"]
    },
    {
      "subsection_title": "Minibatch Construction",
      "explanation": "To train efficiently on massive graphs, PinSage uses a producer-consumer minibatch approach, conceptually illustrated in **diagram img_7_3.png**. This allows parallel preparation of computation graphs on the CPU while the GPU performs model training, maximizing hardware utilization and enabling large batch sizes.",
      "associated_image_filenames": ["img_7_3.png"]
    }
  ]
}
```

**Input Data:**

**Document Text:**
---
""",
        text_content[:30000], # Limit text length
        "\n---\n\n**Available (Pre-filtered) Image Filenames (corresponding to the images provided):**\n",
        json.dumps(image_filenames_list),
        "\n\n**JSON Output:**\n"
    ]

    for img in prepared_images:
        prompt_parts.append(img)

    # Use explicit JSON mode with gemini-1.5-pro
    generation_config = genai.types.GenerationConfig(response_mime_type="application/json")

    try:
        logger.info(f"Sending hierarchical multimodal request to Gemini with {len(prepared_images)} images...")
        response = await gemini_model.generate_content_async(
            prompt_parts,
            generation_config=generation_config
        )

        cleaned_response_text = response.text.strip().strip('```json').strip('```').strip()
        logger.info(f"Received raw response from Gemini: {response.text}")
        logger.info(f"Cleaned response text for JSON parsing: {cleaned_response_text}")

        # Use Pydantic to parse and validate the list of sections
        # Need a temporary wrapper for the top-level list
        class TempWrapper(models.BaseModel):
             # Use the correct hierarchical model name: StudyGuideSection
             study_guide: List[models.StudyGuideSection]

        parsed_wrapper = TempWrapper.model_validate_json(f'{{"study_guide": {cleaned_response_text}}}')
        parsed_sections = parsed_wrapper.study_guide

        logger.info(f"Successfully parsed Gemini response into {len(parsed_sections)} main sections.")
        return parsed_sections

    except Exception as e:
        logger.error(f"Error generating hierarchical study guide or parsing response: {e}")
        logger.error(f"Cleaned text that failed JSON parsing: {cleaned_response_text if 'cleaned_response_text' in locals() else 'N/A'}")
        logger.error(f"Gemini raw response was: {response.text if 'response' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"Failed to generate hierarchical study guide using AI or parse its response: {e}")


# --- API Endpoint ---
# Update the route to the desired structure
@router.post("/workspaces/{workspace_id}/upload", response_model=models.StudyGuideResponse)
async def upload_and_create_study_guide( # Renamed function for clarity
    workspace_id: str, # Add workspace_id path parameter
    file: UploadFile = File(...),
    db: AsyncIOMotorDatabase = Depends(get_database),
    fs: AsyncIOMotorGridFSBucket = Depends(get_gridfs_bucket)
):
    """
    Uploads a document for a specific workspace, extracts text/images,
    and generates a hierarchical study guide using Gemini multimodal capabilities.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    request_id = str(uuid.uuid4())
    request_temp_dir = os.path.join(TEMP_UPLOAD_DIR, request_id)
    request_image_dir = os.path.join(IMAGE_OUTPUT_DIR_BASE, request_id)
    os.makedirs(request_temp_dir, exist_ok=True)
    os.makedirs(request_image_dir, exist_ok=True)

    _, extension = os.path.splitext(file.filename)
    extension = extension.lower()
    temp_file_path = os.path.join(request_temp_dir, f"upload{extension}")

    try:
        logger.info(f"Receiving file: {file.filename}, size: {file.size}, type: {file.content_type}")
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        logger.info(f"Saved temporary file: {temp_file_path}")
    except Exception as e:
        logger.error(f"Failed to save temporary file {temp_file_path}: {e}")
        shutil.rmtree(request_temp_dir, ignore_errors=True)
        shutil.rmtree(request_image_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Could not save uploaded file: {e}")
    finally:
        await file.close()

    text_content = ""
    # Keep track of *all* extracted images before filtering for the final response
    all_extracted_images_info: List[models.ExtractedImageInfo] = []
    # Keep track of images that pass filtering to send to Gemini
    filtered_images_info: List[models.ExtractedImageInfo] = []
    try:
        if extension == ".pdf":
            # Update to capture both lists from extraction function
            text_content, all_extracted_images_info, filtered_images_info = extract_text_and_images_pdf(temp_file_path, request_image_dir)
        elif extension == ".pptx":
             # Update to capture both lists from extraction function
            text_content, all_extracted_images_info, filtered_images_info = extract_text_and_images_pptx(temp_file_path, request_image_dir)
        elif extension == ".docx":
            text_content = extract_text_docx(temp_file_path)
            # No images extracted for docx yet
        elif extension == ".txt":
            text_content = extract_text_txt(temp_file_path)
            # No images for txt
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {extension}. Supported types: .pdf, .pptx, .docx, .txt")

        if not text_content.strip():
            raise HTTPException(status_code=400, detail="Extracted text content is empty.")

        # Generate study guide sections using Gemini with filtered images
        # Use the correct function name and expect hierarchical structure
        study_guide_main_sections = await generate_hierarchical_study_guide(text_content, filtered_images_info, request_image_dir)

        # Format the final response using the correct model
        response_data = models.StudyGuideResponse(
            original_filename=file.filename,
            extracted_images=all_extracted_images_info, # Return info for all originally extracted images
            study_guide=study_guide_main_sections, # Hierarchical study guide
            workspace_id=ObjectId(workspace_id) # Add workspace_id to the model data
        )

        # Construct paths for the filtered images that were actually processed and potentially sent to Gemini
        filtered_image_paths = [os.path.join(request_image_dir, img_info.filename) for img_info in filtered_images_info]

        # Save the generated study guide metadata and upload files to GridFS
        saved_study_guide = await crud.create_study_guide(
            db=db,
            fs=fs, # Pass GridFS bucket
            study_guide_data=response_data,
            original_pdf_path=temp_file_path, # Pass path to original PDF
            image_paths=filtered_image_paths # Pass paths of filtered images
        )

        # Return the saved data (which includes the _id and GridFS IDs)
        return saved_study_guide

    except HTTPException as http_exc:
        raise http_exc
    except Exception as e:
        logger.error(f"Unexpected error processing file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {e}")
    finally:
        shutil.rmtree(request_temp_dir, ignore_errors=True)
        # Only remove image dir if it exists and is not needed anymore
        if os.path.exists(request_image_dir):
             shutil.rmtree(request_image_dir, ignore_errors=True)
        logger.info(f"Cleaned up temporary directories for request {request_id}")

# Endpoint to retrieve a study guide by its ID
@router.get("/{study_guide_id}", response_model=models.StudyGuideResponse)
async def read_study_guide(study_guide_id: str, db: AsyncIOMotorDatabase = Depends(get_database)):
    """
    Retrieves a study guide document by its MongoDB _id.
    Note: This returns the metadata including GridFS IDs, not the file content itself.
    """
    study_guide = await crud.get_study_guide(db=db, study_guide_id=study_guide_id)
    if study_guide is None:
        raise HTTPException(status_code=404, detail="Study guide not found")
    return study_guide

# Endpoint to retrieve all study guides for a specific workspace
@router.get("/workspaces/{workspace_id}/study-guides/", response_model=List[models.StudyGuideResponse]) # Corrected path
async def list_study_guides_for_workspace(workspace_id: str, db: AsyncIOMotorDatabase = Depends(get_database)):
    """
    Retrieves all study guide documents associated with a specific workspace_id.
    """
    # Basic validation for workspace_id format
    if not ObjectId.is_valid(workspace_id):
         raise HTTPException(status_code=400, detail="Invalid workspace ID format")

    study_guides = await crud.get_study_guides_for_workspace(db=db, workspace_id=workspace_id)
    # The CRUD function returns an empty list if none are found or on error,
    # so we don't strictly need to check for None here unless we want to differentiate.
    return study_guides
