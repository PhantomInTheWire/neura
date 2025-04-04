import os
import shutil
import uuid
import json
from fastapi import APIRouter, UploadFile, File, HTTPException, Depends
from typing import List, Dict, Any
import google.generativeai as genai
from dotenv import load_dotenv
import logging
from PIL import Image # For image handling
import io # For handling image bytes

# Text/Image extraction libraries
import fitz # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE # To identify picture shapes in PPTX
from docx import Document
# python-docx does not directly extract images easily, need python-docx-image for that
# Let's add python-docx-image to requirements later if needed, or skip docx images for now.
# For initial implementation, we'll focus on PDF and PPTX images.

from python_chatedu_api import models

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()

# --- Configuration ---
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
TEMP_UPLOAD_DIR = "temp_uploads"
# Create a subdirectory within temp_uploads for each request's images
IMAGE_OUTPUT_DIR_BASE = os.path.join(TEMP_UPLOAD_DIR, "images")

# Configure Gemini
if not GEMINI_API_KEY:
    logger.warning("GEMINI_API_KEY not found. Study guide generation will fail.")
    gemini_model = None
else:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        # Using gemini-1.5-pro which supports multimodal and JSON output mode
        gemini_model = genai.GenerativeModel('gemini-2.0-flash')
        logger.info("Gemini configured successfully with gemini-1.5-pro.")
    except Exception as e:
        logger.error(f"Failed to configure Gemini: {e}")
        gemini_model = None

# Ensure temp directories exist
os.makedirs(TEMP_UPLOAD_DIR, exist_ok=True)
os.makedirs(IMAGE_OUTPUT_DIR_BASE, exist_ok=True)

# --- Helper Functions ---

def extract_text_and_images_pdf(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo]]:
    """Extracts text and images from a PDF file using PyMuPDF."""
    text = ""
    images_info: List[models.ExtractedImageInfo] = []
    try:
        doc = fitz.open(file_path)
        img_counter = 0
        for page_num, page in enumerate(doc):
            text += page.get_text() + "\n"
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                # Generate a unique filename for the image
                image_filename = f"img_{page_num + 1}_{img_counter}.{image_ext}"
                image_save_path = os.path.join(request_image_dir, image_filename)
                with open(image_save_path, "wb") as img_file:
                    img_file.write(image_bytes)
                images_info.append(models.ExtractedImageInfo(filename=image_filename, page_number=page_num + 1))
                img_counter += 1
        logger.info(f"Extracted text and {len(images_info)} images from PDF: {file_path}")
        doc.close()
        return text, images_info
    except Exception as e:
        logger.error(f"Error extracting data from PDF {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PDF file: {e}")

def extract_text_and_images_pptx(file_path: str, request_image_dir: str) -> tuple[str, List[models.ExtractedImageInfo]]:
    """Extracts text and images from a PPTX file."""
    text = ""
    images_info: List[models.ExtractedImageInfo] = []
    try:
        prs = Presentation(file_path)
        img_counter = 0
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                # Check if the shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext.lower()
                    # Generate a unique filename
                    image_filename = f"img_slide_{slide_num + 1}_{img_counter}.{image_ext}"
                    image_save_path = os.path.join(request_image_dir, image_filename)
                    with open(image_save_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    # Page number isn't directly applicable, use slide number
                    images_info.append(models.ExtractedImageInfo(filename=image_filename, page_number=slide_num + 1))
                    img_counter += 1
        logger.info(f"Extracted text and {len(images_info)} images from PPTX: {file_path}")
        return text, images_info
    except Exception as e:
        logger.error(f"Error extracting data from PPTX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX file: {e}")

def extract_text_docx(file_path: str) -> str:
    """Extracts text from a DOCX file. Image extraction is complex and omitted for now."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        logger.info(f"Extracted text from DOCX: {file_path}. Image extraction skipped.")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process DOCX file: {e}")

def extract_text_txt(file_path: str) -> str:
    """Extracts text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        logger.info(f"Extracted text from TXT: {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error extracting text from TXT {file_path}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process TXT file: {e}")

async def generate_enhanced_study_guide(text_content: str, images_info: List[models.ExtractedImageInfo], request_image_dir: str) -> List[models.EnhancedStudyGuideSection]:
    """Uses Gemini multimodal capabilities to generate sections, explanations, and image associations."""
    if not gemini_model:
        raise HTTPException(status_code=500, detail="Gemini model not configured. Check API key.")

    # Prepare image data for Gemini API
    # The exact format depends on the google-generativeai library version
    # Usually involves creating Image objects or similar structures
    prepared_images = []
    image_filenames_list = []
    for img_info in images_info:
        try:
            img_path = os.path.join(request_image_dir, img_info.filename)
            # Example using Pillow to load image data - adjust based on genai library requirements
            pil_image = Image.open(img_path)
            # Convert to RGB if necessary (some models might require it)
            if pil_image.mode != 'RGB':
                 pil_image = pil_image.convert('RGB')
            prepared_images.append(pil_image) # The library might handle PIL images directly
            image_filenames_list.append(img_info.filename)
        except Exception as e:
            logger.warning(f"Could not load or prepare image {img_info.filename}: {e}")

    # Refined prompt v2: Detailed explanation and image integration
    prompt_parts = [
        """You are an expert academic assistant tasked with creating an enhanced, explanatory study guide from a document's text and associated images. Your goal is to break down the content into logical sections, **thoroughly explain** the concepts within each section like a tutor would, and **integrate explanations of relevant images** directly into the text.

**Input:**
You will receive the full text content of the document and a list of image filenames corresponding to images extracted from the document that will be provided alongside the text.

**Tasks:**

1.  **Analyze Text:** Read through the entire provided text content.
2.  **Identify Sections:** Divide the text into logical sections based on the main topics or concepts discussed.
3.  **Generate Explanations (Detailed & Integrated):** For each identified section:
    *   Create a concise and descriptive `section_title`.
    *   Write a detailed `explanation` (aim for **at least 5-8 sentences**, or more if needed for clarity) that:
        *   Thoroughly clarifies the key concepts discussed within that section's text.
        *   Defines important jargon or technical terms.
        *   Simplifies complex ideas using analogies or simpler language where appropriate.
        *   Provides necessary background context if helpful for understanding.
        *   **Use Markdown syntax within this explanation text for formatting:** Use asterisks (`*` or `-`) for bullet points, double asterisks (`**text**`) for bolding, backticks (`) for inline code, and LaTeX delimiters (`$...$` or `$$...$$`) for mathematical formulas. Use Markdown table syntax if needed.
        *   **Crucially:** If you associate any images with this section (see Task 4), **explicitly reference them within this explanation text** (e.g., "Figure [filename] illustrates...", "As seen in the diagram [filename], the process involves...", "The results in [filename] show..."). Explain what the referenced image depicts and how it connects to the concepts being explained.
    *   **Focus:** Prioritize deep explanation and enhancing understanding over mere summarization. Do not simply rephrase the original text; add value through clarification and context.
4.  **Associate Images:** For each section, review the provided images (which correspond to the filenames listed below) and determine which images are directly relevant to the content being explained in that section.
5.  **Format Output:** Structure your entire response as a single JSON list of objects. Each object represents a section and must contain the following keys:
    *   `section_title` (string): The title you generated for the section.
    *   `explanation` (string): The detailed explanatory text you generated, including references to associated images.
    *   `associated_image_filenames` (list of strings): A list containing the filenames of the images you identified as relevant *and referenced* in the explanation. If no images are relevant or referenced, provide an empty list `[]`.

**Output Constraints:**
*   Your response MUST be only the JSON list, starting with `[` and ending with `]`.
*   Do NOT include any introductory text, concluding remarks, or markdown formatting (like ```json) around the JSON output.
*   Ensure the generated JSON is valid.

**Example JSON Object within the List:**
```json
{
  "section_title": "PinSage Architecture Overview",
  "explanation": "The PinSage model utilizes localized convolutions, as depicted in Figure img_6_2.png. This means that instead of processing the entire graph at once, it focuses on a node's immediate neighborhood. The process involves sampling neighbors using random walks and then aggregating their features, weighted by importance. This localized approach, combined with efficient minibatching shown conceptually in img_7_3.png, allows the model to scale effectively to billions of nodes without requiring the full graph Laplacian, overcoming limitations of earlier GCNs.",
  "associated_image_filenames": ["img_6_2.png", "img_7_3.png"]
}
```

**Input Data:**

**Document Text:**
---
""",
        text_content[:30000], # Limit text length if needed for API/model
        "\n---\n\n**Available Image Filenames (corresponding to the images provided in the input):**\n",
        json.dumps(image_filenames_list), # Send filenames for reference in the prompt
        "\n\n**JSON Output:**\n"
    ]

    # Add prepared images to the prompt parts list
    # The exact method depends on the library version (e.g., appending PIL Images)
    for img in prepared_images:
        prompt_parts.append(img)

    # generation_config removed as the experimental model doesn't support it.
    # Relying on prompt instructions AND explicit JSON mode for format.
    try:
        logger.info(f"Sending multimodal request to Gemini with {len(prepared_images)} images...")
        # Re-adding generation_config for explicit JSON output mode
        response = await gemini_model.generate_content_async(
            prompt_parts,
            generation_config={"response_mime_type": "application/json"} # Pass dict directly
        )

        # Attempt to parse the response as JSON
        # With response_mime_type set, cleaning might be less necessary, but keep basic strip just in case
        cleaned_response_text = response.text.strip().strip('```json').strip('```').strip()
        logger.info(f"Received raw response from Gemini: {response.text}")
        logger.info(f"Cleaned response text for JSON parsing: {cleaned_response_text}")

        # Use Pydantic to parse and validate the list of sections
        # Need a temporary wrapper structure for validation if the response is just the list
        class TempWrapper(models.BaseModel):
             study_guide: List[models.EnhancedStudyGuideSection]

        parsed_wrapper = TempWrapper.model_validate_json(f'{{"study_guide": {cleaned_response_text}}}')
        parsed_sections = parsed_wrapper.study_guide

        logger.info(f"Successfully parsed Gemini response into {len(parsed_sections)} sections.")
        return parsed_sections

    except Exception as e:
        logger.error(f"Error generating study guide with Gemini or parsing response: {e}")
        # Log the cleaned text that failed parsing
        logger.error(f"Cleaned text that failed JSON parsing: {cleaned_response_text if 'cleaned_response_text' in locals() else 'N/A'}")
        # Log the raw response if available
        logger.error(f"Gemini raw response was: {response.text if 'response' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"Failed to generate study guide using AI or parse its response: {e}")


# --- API Endpoint ---

@router.post("/upload", response_model=models.EnhancedStudyGuideResponse)
async def create_enhanced_study_guide(file: UploadFile = File(...)):
    """
    Uploads a document (PDF, PPTX, DOCX, TXT), extracts text and images (PDF, PPTX),
    and generates an enhanced study guide using Gemini multimodal capabilities.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    # Create a unique directory for this request's temporary files
    request_id = str(uuid.uuid4())
    request_temp_dir = os.path.join(TEMP_UPLOAD_DIR, request_id)
    request_image_dir = os.path.join(IMAGE_OUTPUT_DIR_BASE, request_id)
    os.makedirs(request_temp_dir, exist_ok=True)
    os.makedirs(request_image_dir, exist_ok=True)

    _, extension = os.path.splitext(file.filename)
    extension = extension.lower()
    temp_file_path = os.path.join(request_temp_dir, f"upload{extension}")

    # Save uploaded file temporarily
    try:
        logger.info(f"Receiving file: {file.filename}, size: {file.size}, type: {file.content_type}")
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        logger.info(f"Saved temporary file: {temp_file_path}")
    except Exception as e:
        logger.error(f"Failed to save temporary file {temp_file_path}: {e}")
        # Clean up directories if save fails
        shutil.rmtree(request_temp_dir, ignore_errors=True)
        shutil.rmtree(request_image_dir, ignore_errors=True)
        raise HTTPException(status_code=500, detail=f"Could not save uploaded file: {e}")
    finally:
        await file.close()

    # Extract text and images based on file type
    text_content = ""
    images_info: List[models.ExtractedImageInfo] = []
    try:
        if extension == ".pdf":
            text_content, images_info = extract_text_and_images_pdf(temp_file_path, request_image_dir)
        elif extension == ".pptx":
            text_content, images_info = extract_text_and_images_pptx(temp_file_path, request_image_dir)
        elif extension == ".docx":
            # Currently only extracts text for docx
            text_content = extract_text_docx(temp_file_path)
            images_info = [] # No images extracted for docx yet
        elif extension == ".txt":
            text_content = extract_text_txt(temp_file_path)
            images_info = [] # No images for txt
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {extension}. Supported types: .pdf, .pptx, .docx, .txt")

        if not text_content.strip():
            raise HTTPException(status_code=400, detail="Extracted text content is empty.")

        # Generate study guide sections using Gemini
        study_guide_sections = await generate_enhanced_study_guide(text_content, images_info, request_image_dir)

        # Format the final response
        response_data = models.EnhancedStudyGuideResponse(
            original_filename=file.filename,
            extracted_images=images_info, # List of all extracted images
            study_guide=study_guide_sections # Sections with explanations and associated image filenames
        )
        return response_data

    except HTTPException as http_exc:
        # Re-raise HTTPExceptions directly
        raise http_exc
    except Exception as e:
        # Catch any other unexpected errors during processing
        logger.error(f"Unexpected error processing file {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {e}")
    finally:
        # Clean up the temporary file and image directories for this request
        shutil.rmtree(request_temp_dir, ignore_errors=True)
        shutil.rmtree(request_image_dir, ignore_errors=True)
        logger.info(f"Cleaned up temporary directories for request {request_id}")
